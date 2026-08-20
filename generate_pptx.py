import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # blank layout

    # Color Palette
    BG_DARK = RGBColor(15, 23, 42)      # Slate 900
    BG_LIGHT = RGBColor(248, 250, 252)  # Slate 50
    CARD_BG = RGBColor(255, 255, 255)   # Pure White
    CARD_DARK = RGBColor(30, 41, 59)    # Slate 800
    BORDER_COLOR = RGBColor(226, 232, 240) # Slate 200
    
    PRIMARY = RGBColor(37, 99, 235)     # Blue 600
    PRIMARY_LIGHT = RGBColor(239, 246, 255) # Blue 50
    ACCENT = RGBColor(16, 185, 129)     # Emerald 500
    ACCENT_LIGHT = RGBColor(236, 253, 245) # Emerald 50
    AMBER = RGBColor(245, 158, 11)      # Amber 500
    AMBER_LIGHT = RGBColor(254, 243, 199) # Amber 100
    PURPLE = RGBColor(139, 92, 246)     # Purple 500
    PURPLE_LIGHT = RGBColor(245, 243, 255)
    
    TEXT_DARK = RGBColor(15, 23, 42)    # Slate 900
    TEXT_MUTED = RGBColor(100, 116, 139) # Slate 500
    TEXT_LIGHT = RGBColor(248, 250, 252) # Slate 50

    def add_header(slide, title_text, category_text="SMART INDIA HACKATHON"):
        # Header Badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.45), Inches(3.2), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = PRIMARY_LIGHT
        badge.line.color.rgb = PRIMARY
        badge.line.width = Pt(1)
        tf = badge.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        p.alignment = PP_ALIGN.CENTER

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.7), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK

    def set_slide_background(slide, color):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=BORDER_COLOR):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
        return card

    # ==========================================
    # SLIDE 1: Title Slide (Dark Theme)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, BG_DARK)

    # Accent Top Line
    top_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = PRIMARY
    top_bar.line.fill.background()

    # Hackathon Badge
    badge1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.2), Inches(4.5), Inches(0.45))
    badge1.fill.solid()
    badge1.fill.fore_color.rgb = RGBColor(30, 58, 138)
    badge1.line.color.rgb = PRIMARY
    tf1 = badge1.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = "🏆 SMART INDIA HACKATHON (SIH) PRESENTATION"
    p1.font.size = Pt(12)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(147, 197, 253)
    p1.alignment = PP_ALIGN.CENTER

    # Main Title
    tbox = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(2.2))
    tf = tbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Skill-Sync AI"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    p_sub = tf.add_paragraph()
    p_sub.text = "Smart Allocation Engine & Career Enablement for PM Internship Scheme"
    p_sub.font.size = Pt(24)
    p_sub.font.bold = True
    p_sub.font.color.rgb = RGBColor(96, 165, 250)

    p_desc = tf.add_paragraph()
    p_desc.text = "An intelligent, multi-tenant platform ensuring merit-driven, equitable, and AI-optimized internship matching at national scale."
    p_desc.font.size = Pt(16)
    p_desc.font.color.rgb = RGBColor(203, 213, 225)

    # 3 Key Value Props Card Row
    c1 = add_card(slide1, 1.0, 4.3, 3.5, 2.2, CARD_DARK, RGBColor(51, 65, 85))
    tf = c1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🎯 Hybrid AI Matching"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(147, 197, 253)
    p2 = tf.add_paragraph()
    p2.text = "\nCombines NLP semantic vector similarity + Jaccard skill coverage + non-linear boosting for optimal student-role fit."
    p2.font.size = Pt(13)
    p2.font.color.rgb = RGBColor(226, 232, 240)

    c2 = add_card(slide1, 4.9, 4.3, 3.5, 2.2, CARD_DARK, RGBColor(51, 65, 85))
    tf = c2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "📊 Diversity & Affirmative Action"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(110, 231, 183)
    p2 = tf.add_paragraph()
    p2.text = "\nReal-time analytics & macro-allocation balancing Gender Parity, Rural/Urban ratios, and Tier-1/2/3 College representation."
    p2.font.size = Pt(13)
    p2.font.color.rgb = RGBColor(226, 232, 240)

    c3 = add_card(slide1, 8.8, 4.3, 3.5, 2.2, CARD_DARK, RGBColor(51, 65, 85))
    tf = c3.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🎙️ AI Upskilling & Mock Studio"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(252, 211, 77)
    p2 = tf.add_paragraph()
    p2.text = "\nAutomated resume parsing, skill-gap diagnostics, 3-phase customized roadmaps, and live AI interview assessment."
    p2.font.size = Pt(13)
    p2.font.color.rgb = RGBColor(226, 232, 240)

    # ==========================================
    # SLIDE 2: Problem Statement & National Context
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, BG_LIGHT)
    add_header(slide2, "The Problem: Scaling the PM Internship Scheme", "PROBLEM & CONTEXT")

    # Left Box: The Mandate
    c_mandate = add_card(slide2, 0.8, 1.8, 4.2, 5.0, PRIMARY_LIGHT, PRIMARY)
    tf = c_mandate.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🇮🇳 The PM Scheme Vision"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    bullets = [
        ("1 Crore Youth", "Targeted over 5 years across India's Top 500 Companies."),
        ("12-Month Internships", "Bridging academic education and industry-grade job readiness."),
        ("Rs 5,000 Monthly Stipend", "Empowering students from all socio-economic backgrounds.")
    ]
    for title, desc in bullets:
        p_t = tf.add_paragraph()
        p_t.text = f"\n• {title}"
        p_t.font.bold = True
        p_t.font.size = Pt(15)
        p_t.font.color.rgb = TEXT_DARK
        p_d = tf.add_paragraph()
        p_d.text = f"   {desc}"
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = TEXT_MUTED

    # Right Column: 4 Pain Points Cards
    pain_points = [
        ("🔴 Massive Scale & Manual Bottlenecks", "Processing millions of diverse resumes against thousands of job descriptions leads to misallocation, delays, and screening fatigue.", 1.8),
        ("🔴 Skill-Opportunity Asymmetry", "Candidates lack visibility into exact industry expectations, while recruiters struggle to evaluate real competency from text resumes.", 3.0),
        ("🔴 Socio-Demographic Inequity", "Students from rural areas, Tier-2/3 institutions, and underrepresented genders face structural disadvantages in discovery.", 4.2),
        ("🔴 Lack of Pre-Placement Enablement", "Applicants from non-metro areas lack access to personalized resume reviews, interview coaching, and skill-gap roadmaps.", 5.4),
    ]

    for title, desc, top_y in pain_points:
        c = add_card(slide2, 5.3, top_y, 7.2, 1.05, CARD_BG, BORDER_COLOR)
        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(220, 38, 38)
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = TEXT_DARK

    # ==========================================
    # SLIDE 3: The Skill-Sync AI Solution
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, BG_LIGHT)
    add_header(slide3, "Skill-Sync AI: A Tri-Party Intelligent Ecosystem", "SOLUTION OVERVIEW")

    roles = [
        ("🎓 Candidate / Student Portal", PRIMARY, PRIMARY_LIGHT, [
            "AI Resume Auto-Parser & Profile Extractor",
            "Skill Gap Matrix with Priority & Learning Time",
            "3-Phase Personalized Upskilling Roadmap",
            "Interactive AI Mock Interviewer with instant scorecard",
            "Continuous Mini-Project Skill Simulator"
        ], 0.8),
        ("🏢 Corporate Recruiter Cockpit", ACCENT, ACCENT_LIGHT, [
            "Structured Internship Requisition Publishing",
            "AI Match Scoreboard & Talent Pipeline",
            "Multi-Criteria Candidate Dossier Inspection",
            "Frictionless Status Progression & Shortlisting",
            "Targeted Skill Requirements Definition"
        ], 4.9),
        ("🏛️ Scheme Governance / Admin", PURPLE, PURPLE_LIGHT, [
            "Automated Global Macro-Allocation Routine",
            "Real-time Diversity & Equity Dashboard",
            "Institutional Tier & Gender Parity Monitoring",
            "Regional & Sectoral Coverage Analytics",
            "User Access Control & Audit Governance"
        ], 9.0)
    ]

    for title, color, bg_c, features, left_x in roles:
        c = add_card(slide3, left_x, 1.8, 3.55, 5.0, CARD_BG, color)
        # Header banner inside card
        banner = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_x + 0.15), Inches(1.95), Inches(3.25), Inches(0.65))
        banner.fill.solid()
        banner.fill.fore_color.rgb = bg_c
        banner.line.color.rgb = color
        tf_b = banner.text_frame
        p_b = tf_b.paragraphs[0]
        p_b.text = title
        p_b.font.size = Pt(13)
        p_b.font.bold = True
        p_b.font.color.rgb = color
        p_b.alignment = PP_ALIGN.CENTER

        tf_c = c.text_frame
        tf_c.word_wrap = True
        p_init = tf_c.paragraphs[0]
        p_init.text = "\n\n"
        for feat in features:
            p_f = tf_c.add_paragraph()
            p_f.text = f"✓  {feat}"
            p_f.font.size = Pt(12)
            p_f.font.color.rgb = TEXT_DARK
            p_f.space_after = Pt(8)

    # ==========================================
    # SLIDE 4: AI Matching Engine & Mathematical Model
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, BG_LIGHT)
    add_header(slide4, "Proprietary AI Matching Algorithm & Mathematical Formulation", "MATCHING ALGORITHM")

    # Left Box: Mathematical Formulation
    c_math = add_card(slide4, 0.8, 1.8, 5.6, 5.0, CARD_BG, PRIMARY)
    tf = c_math.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "📐 Hybrid Matching Formulation"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    p_eq1 = tf.add_paragraph()
    p_eq1.text = "\n1. Raw Multi-Factor Score:"
    p_eq1.font.bold = True
    p_eq1.font.size = Pt(14)
    p_eq1.font.color.rgb = TEXT_DARK

    p_eq2 = tf.add_paragraph()
    p_eq2.text = "   Score_raw = (0.30 × S_vector) + (0.70 × S_skill)"
    p_eq2.font.bold = True
    p_eq2.font.size = Pt(14)
    p_eq2.font.color.rgb = RGBColor(37, 99, 235)

    p_eq3 = tf.add_paragraph()
    p_eq3.text = "\n2. Non-Linear Score Boosting Curve:"
    p_eq3.font.bold = True
    p_eq3.font.size = Pt(14)
    p_eq3.font.color.rgb = TEXT_DARK

    p_eq4 = tf.add_paragraph()
    p_eq4.text = "   Score_final = (Score_raw)^0.5 × 100"
    p_eq4.font.bold = True
    p_eq4.font.size = Pt(14)
    p_eq4.font.color.rgb = RGBColor(16, 185, 129)

    p_desc = tf.add_paragraph()
    p_desc.text = "\nWhy this works:\n• 70% weight on mandatory skills guarantees role competency.\n• 30% semantic NLP captures domain aspirations & soft skills.\n• Square-root curve prevents harsh penalties for candidates with solid foundations, encouraging continuous upskilling."
    p_desc.font.size = Pt(12)
    p_desc.font.color.rgb = TEXT_MUTED

    # Right Cards: Algorithm Steps
    steps = [
        ("Step 1: NLP Vectorization", "Vectorizes student career goals, skills, and industry preferences alongside role descriptions using TF-IDF with English stop-word filtration.", PRIMARY),
        ("Step 2: Semantic Cosine Similarity", "Computes high-dimensional cosine angle S_vector between candidate profile vectors and corporate requisition embeddings.", ACCENT),
        ("Step 3: Jaccard Skill Overlap", "Calculates exact set intersection: S_skill = |Student Skills ∩ Required Skills| / |Required Skills|.", AMBER),
        ("Step 4: Real-time Multi-Candidate Ranking", "Ranks top candidates and displays transparent scorecards with drill-down breakdown to candidates and recruiters.", PURPLE)
    ]

    for idx, (st_title, st_desc, color) in enumerate(steps):
        top_y = 1.8 + (idx * 1.25)
        c = add_card(slide4, 6.7, top_y, 5.8, 1.15, CARD_BG, BORDER_COLOR)
        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = st_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color
        p_d = tf.add_paragraph()
        p_d.text = st_desc
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = TEXT_DARK

    # ==========================================
    # SLIDE 5: Technical Architecture & Tech Stack
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, BG_LIGHT)
    add_header(slide5, "End-to-End Scalable Technology Architecture", "SYSTEM ARCHITECTURE")

    tech_layers = [
        ("Frontend Application", PRIMARY, [
            ("Framework", "React 19.1 + TypeScript 5.8"),
            ("Styling", "Tailwind CSS + Lucide Icons"),
            ("Build Tool", "Vite 6.2 (Sub-second HMR)"),
            ("State & Routing", "React Router DOM + Context API")
        ], 0.8),
        ("Backend REST Services", ACCENT, [
            ("Core Engine", "Python 3.9+ & FastAPI 0.128"),
            ("Async Server", "Uvicorn & AsyncIO"),
            ("Security", "JWT Tokens & Passlib Bcrypt"),
            ("Architecture", "Modular Endpoints & Pydantic")
        ], 3.8),
        ("AI & Machine Learning", AMBER, [
            ("LLM Inference", "Groq Cloud (Qwen 3.6 / Llama 3)"),
            ("NLP Matching", "Scikit-Learn TF-IDF Vectorizer"),
            ("Similarity", "Cosine Similarity + NumPy"),
            ("Fallback AI", "Local Rule-based Heuristics")
        ], 6.8),
        ("Persistence & Database", PURPLE, [
            ("Database", "MongoDB Atlas / Local NoSQL"),
            ("Async ODM", "Beanie ODM Framework"),
            ("Async Driver", "Motor Python Client"),
            ("Data Schemas", "Pydantic V2 Document Models")
        ], 9.8)
    ]

    for title, color, items, left_x in tech_layers:
        c = add_card(slide5, left_x, 1.8, 2.7, 5.0, CARD_BG, color)
        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        
        for k, v in items:
            p_k = tf.add_paragraph()
            p_k.text = f"\n{k}:"
            p_k.font.bold = True
            p_k.font.size = Pt(12)
            p_k.font.color.rgb = TEXT_DARK
            p_v = tf.add_paragraph()
            p_v.text = v
            p_v.font.size = Pt(12)
            p_v.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 6: Student Upskilling & Career Engine
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, BG_LIGHT)
    add_header(slide6, "Candidate Enablement: Upskilling, Resume AI & Roadmaps", "STUDENT ENABLEMENT")

    features_s6 = [
        ("📄 AI Resume Auto-Extractor", "Upload raw text or resume files. Groq LLM parses candidate skills, projects, degrees, strengths, and areas of improvement into structured JSON schemas automatically.", PRIMARY, 0.8, 1.8),
        ("🔍 Real-Time Skill-Gap Matrix", "Identifies exact missing competencies for target internships. Classifies skills into High/Medium/Low priority with estimated time-to-learn (e.g. '2 weeks').", ACCENT, 6.7, 1.8),
        ("🗺️ 3-Phase Personalized Roadmap", "Generates structured 12-week learning pathways: Phase 1 (Foundations), Phase 2 (Applied Systems), Phase 3 (Portfolio & Interview Readiness) with course links & capstone ideas.", PURPLE, 0.8, 4.3),
        ("⚡ Interactive Project Simulator", "Candidates complete practical mini-projects (e.g. Market Analysis Report). AI grades submissions and instantly recalculates match scores across live listings.", AMBER, 6.7, 4.3),
    ]

    for title, desc, color, left_x, top_y in features_s6:
        c = add_card(slide6, left_x, top_y, 5.8, 2.3, CARD_BG, BORDER_COLOR)
        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        p_d = tf.add_paragraph()
        p_d.text = f"\n{desc}"
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = TEXT_DARK

    # ==========================================
    # SLIDE 7: AI Mock Interview Studio
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, BG_LIGHT)
    add_header(slide7, "AI Mock Interview Studio: Real-Time Performance Analytics", "MOCK INTERVIEW STUDIO")

    # Left: Mock Interview Flow
    c_flow = add_card(slide7, 0.8, 1.8, 5.6, 5.0, CARD_BG, PRIMARY)
    tf = c_flow.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🎙️ Live Simulated Interview Flow"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    flow_items = [
        ("1. Role Selection", "Candidate selects specific target role (e.g. AI Product Analyst, Data Engineer)."),
        ("2. Interactive Video/Voice Session", "AI asks structured behavioral (STAR) and domain technical questions."),
        ("3. Speech & Keyword Analysis", "Analyzes technical terminology coverage, response cadence, and clarity."),
        ("4. Comprehensive Diagnostic Report", "Generates overall score (e.g. 88/100), key strengths, and exact improvement suggestions.")
    ]
    for title, desc in flow_items:
        p_t = tf.add_paragraph()
        p_t.text = f"\n• {title}"
        p_t.font.bold = True
        p_t.font.size = Pt(14)
        p_t.font.color.rgb = TEXT_DARK
        p_d = tf.add_paragraph()
        p_d.text = f"   {desc}"
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = TEXT_MUTED

    # Right: Key Metrics Evaluated
    c_metrics = add_card(slide7, 6.7, 1.8, 5.8, 5.0, CARD_BG, ACCENT)
    tf = c_metrics.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "📊 4-Pillar Evaluation Matrix"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    matrix_items = [
        ("1. Verbal Clarity & Articulation", "Measures conciseness, structured thoughts, and elimination of filler words."),
        ("2. Technical & Domain Keywords", "Evaluates accurate utilization of role-specific tech stack and frameworks."),
        ("3. Behavioral & STAR Structuring", "Checks Situation-Task-Action-Result methodology in answering situational prompts."),
        ("4. Actionable Remediation Plan", "Provides specific links to resources for areas with low readiness scores.")
    ]
    for title, desc in matrix_items:
        p_t = tf.add_paragraph()
        p_t.text = f"\n• {title}"
        p_t.font.bold = True
        p_t.font.size = Pt(14)
        p_t.font.color.rgb = TEXT_DARK
        p_d = tf.add_paragraph()
        p_d.text = f"   {desc}"
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 8: Admin Governance & Diversity Dashboard
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8, BG_LIGHT)
    add_header(slide8, "Scheme Governance & Diversity Analytics Cockpit", "DIVERSITY & GOVERNANCE")

    div_cards = [
        ("⚧️ Gender Parity Monitoring", "Real-time tracking of female and non-binary representation across sectors, ensuring equal opportunity across STEM and business domains.", PRIMARY, 0.8, 1.8),
        ("🏞️ Rural vs. Urban Equity", "Ensures candidates from rural and semi-urban districts are not marginalized by geographical location or access barriers.", ACCENT, 6.7, 1.8),
        ("🏛️ Tier-1 / 2 / 3 College Inclusivity", "Actively monitors college tier representation to ensure merit from Tier-2 and Tier-3 institutions surfaces to top 500 corporate recruiters.", PURPLE, 0.8, 4.3),
        ("⚙️ Automated Macro Allocation Routine", "One-click execution of national-scale batch matching algorithm balancing corporate vacancies, merit scores, and diversity targets.", AMBER, 6.7, 4.3),
    ]

    for title, desc, color, left_x, top_y in div_cards:
        c = add_card(slide8, left_x, top_y, 5.8, 2.3, CARD_BG, BORDER_COLOR)
        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        p_d = tf.add_paragraph()
        p_d.text = f"\n{desc}"
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = TEXT_DARK

    # ==========================================
    # SLIDE 9: Corporate Recruiter Cockpit
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9, BG_LIGHT)
    add_header(slide9, "Corporate Recruiter Portal: Frictionless Talent Pipeline", "CORPORATE WORKFLOW")

    corp_features = [
        ("1. Intuitive Requisition Builder", "Post role details including mandatory skills, sector, location, duration, stipend, and openings in seconds."),
        ("2. AI Ranked Match Leaderboard", "Automatic ranking of candidates by hybrid match score, highlighting top-percentile talent instantly."),
        ("3. Candidate Dossier & Analytics", "View comprehensive candidate profile, verified certificates, mini-project scores, and AI mock interview feedback."),
        ("4. One-Click Pipeline Actions", "Seamlessly shortlist, accept, or reject candidates with automated real-time candidate notifications.")
    ]

    for idx, (title, desc) in enumerate(corp_features):
        top_y = 1.8 + (idx * 1.25)
        c = add_card(slide9, 0.8, top_y, 11.7, 1.15, CARD_BG, BORDER_COLOR)
        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = TEXT_DARK

    # ==========================================
    # SLIDE 10: Live Demonstration Workflow
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10, BG_LIGHT)
    add_header(slide10, "Live Demonstration Walkthrough for Jury", "DEMO WALKTHROUGH")

    demo_steps = [
        ("Step 1: Student Experience", PRIMARY, [
            "Log in as candidate: student@example.com",
            "View hybrid AI match scores for internships",
            "Upload resume & review extracted skills & projects",
            "Inspect Skill Gap matrix & 3-phase roadmap",
            "Complete mini-project simulator & observe score boost"
        ], 0.8),
        ("Step 2: Recruiter Experience", ACCENT, [
            "Log in as company: company@example.com",
            "Publish new internship requisition",
            "Review AI ranked applicant leaderboard",
            "Inspect top applicant's detailed dossier",
            "Shortlist candidate & trigger notification"
        ], 4.9),
        ("Step 3: Admin Experience", PURPLE, [
            "Log in as admin: admin@example.com",
            "Inspect Gender, Rural/Urban & Tier ratios",
            "Trigger Automated Macro Allocation engine",
            "Inspect system telemetry & user governance"
        ], 9.0)
    ]

    for title, color, steps_list, left_x in demo_steps:
        c = add_card(slide10, left_x, 1.8, 3.55, 5.0, CARD_BG, color)
        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        
        for st in steps_list:
            p_s = tf.add_paragraph()
            p_s.text = f"\n• {st}"
            p_s.font.size = Pt(12)
            p_s.font.color.rgb = TEXT_DARK

    # ==========================================
    # SLIDE 11: Scalability, Impact & Future Roadmap
    # ==========================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11, BG_LIGHT)
    add_header(slide11, "National Impact, Scalability & Strategic Roadmap", "FUTURE ROADMAP")

    future_milestones = [
        ("🌐 Multi-Lingual Regional Support", "Integrate vernacular Indian language models (Bhashini AI / IndicLLM) to support rural applicants in Hindi, Tamil, Telugu, Marathi, and Bengali.", PRIMARY, 0.8, 1.8),
        ("🔗 Blockchain Verified Credentials", "Tamper-proof verifiable credentials on DigiLocker / IndiaStack blockchain for automatic verification of degrees and certifications.", ACCENT, 6.7, 1.8),
        ("📱 Native Mobile Application", "PWA & low-bandwidth Flutter app with offline caching and SMS notification fallback for remote areas with intermittent connectivity.", PURPLE, 0.8, 4.3),
        ("🚀 Pan-India Kubernetes Microservices", "Containerized microservice scaling on MeghRaj / NIC Cloud capable of handling 50,000+ concurrent allocation requests per second.", AMBER, 6.7, 4.3)
    ]

    for title, desc, color, left_x, top_y in future_milestones:
        c = add_card(slide11, left_x, top_y, 5.8, 2.3, CARD_BG, BORDER_COLOR)
        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        p_d = tf.add_paragraph()
        p_d.text = f"\n{desc}"
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = TEXT_DARK

    # ==========================================
    # SLIDE 12: Summary & Q&A
    # ==========================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide12, BG_DARK)

    # Accent Top Line
    top_bar = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = ACCENT
    top_bar.line.fill.background()

    # Thank You Title
    tbox = slide12.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.3), Inches(2.0))
    tf = tbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Thank You! | Questions & Answers"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    p_sub = tf.add_paragraph()
    p_sub.text = "Skill-Sync AI: Empowering India's Youth Through Transparent, Equitable & Smart Allocation"
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = RGBColor(147, 197, 253)

    # Summary Card
    c_summary = add_card(slide12, 1.0, 3.4, 11.3, 3.4, CARD_DARK, RGBColor(51, 65, 85))
    tf_s = c_summary.text_frame
    tf_s.word_wrap = True
    p_s = tf_s.paragraphs[0]
    p_s.text = "🔑 Why Skill-Sync AI Wins for PM Internship Scheme:"
    p_s.font.size = Pt(18)
    p_s.font.bold = True
    p_s.font.color.rgb = RGBColor(252, 211, 77)

    summary_bullets = [
        "1. Complete End-to-End Working System (React 19 + FastAPI + MongoDB + Groq AI).",
        "2. Mathematically Sound Hybrid Matching (TF-IDF + Jaccard + Non-linear Boosting).",
        "3. Solves Real National Bottlenecks: Inclusivity, Skill Gaps, and Pre-Placement Training.",
        "4. Live Interactive AI Capabilities: Mock Interviews, Resume Parser, and 3-Phase Roadmaps.",
        "5. Ready for National Deployment across India's top 500 companies."
    ]
    for b in summary_bullets:
        p_b = tf_s.add_paragraph()
        p_b.text = f"•  {b}"
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = RGBColor(226, 232, 240)
        p_b.space_after = Pt(4)

    output_path = os.path.abspath("SIH_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_deck()
